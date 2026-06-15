import os
import re
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

def set_fade_transition(slide):
    transition = OxmlElement('p:transition')
    fade = OxmlElement('p:fade')
    transition.append(fade)
    slide._element.insert(2, transition)

def clean_text(text):
    return re.sub(r'\s+', ' ', text).strip()

def process_html_to_pptx(base_dir, output_path):
    files = [
        "index.html",
        "chapter-01.html",
        "chapter-02.html",
        "chapter-03.html",
        "chapter-04.html",
        "chapter-05.html",
        "chapter-06.html",
        "chapter-07.html",
        "appendices.html",
        "bibliography.html",
        "academic-engagements.html"
    ]
    
    prs = Presentation()
    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6]
    
    for file_name in files:
        file_path = os.path.join(base_dir, file_name)
        if not os.path.exists(file_path):
            print(f"Skipping {file_name}, not found.")
            continue
            
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            
        slides = soup.find_all('section', class_='slide')
        
        for slide_html in slides:
            slide = prs.slides.add_slide(blank_slide_layout)
            try:
                set_fade_transition(slide)
            except Exception as e:
                pass
                
            current_top = Inches(0.5)
            left_margin = Inches(0.5)
            slide_width = prs.slide_width
            max_width = slide_width - Inches(1)
            
            # Title
            title_elem = slide_html.find(class_='slide-title')
            if not title_elem:
                title_elem = slide_html.find(['h1', 'h2'])
            
            if title_elem:
                title_text = clean_text(title_elem.get_text())
                txBox = slide.shapes.add_textbox(left_margin, current_top, max_width, Inches(1))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.name = 'Arial'
                p.font.color.rgb = RGBColor(0, 0, 0)
                current_top += Inches(1.0)
                # Remove from tree so we don't process it again
                title_elem.extract()
                
            # Subtitle
            subtitle_elem = slide_html.find(class_='slide-subtitle')
            if subtitle_elem:
                subtitle_text = clean_text(subtitle_elem.get_text())
                txBox = slide.shapes.add_textbox(left_margin, current_top, max_width, Inches(0.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = subtitle_text
                p.font.size = Pt(24)
                p.font.name = 'Arial'
                p.font.color.rgb = RGBColor(100, 100, 100)
                current_top += Inches(0.6)
                subtitle_elem.extract()
                
            # Now let's try to process text, images, and tables in order of appearance
            # We will flatten the structure somewhat
            for element in slide_html.find_all(['h3', 'h4', 'p', 'ul', 'ol', 'table', 'img', 'div']):
                # skip elements inside tables or uls if we process the parent
                if element.find_parent('table') and element.name != 'table':
                    continue
                if element.find_parent('ul') and element.name != 'ul':
                    continue
                if element.find_parent('ol') and element.name != 'ol':
                    continue
                if element.name == 'div' and not ('glass-card' in element.get('class', []) or 'stat-item' in element.get('class', [])):
                    continue
                    
                if element.name in ['h3', 'h4']:
                    text = clean_text(element.get_text())
                    if not text: continue
                    txBox = slide.shapes.add_textbox(left_margin, current_top, max_width, Inches(0.4))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    current_top += Inches(0.5)
                    element.extract()
                    
                elif element.name == 'p':
                    text = clean_text(element.get_text())
                    if not text: continue
                    txBox = slide.shapes.add_textbox(left_margin, current_top, max_width, Inches(0.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(50, 50, 50)
                    # estimate height
                    current_top += Inches(0.5 + (len(text) // 100) * 0.2)
                    element.extract()
                    
                elif element.name in ['ul', 'ol']:
                    items = [clean_text(li.get_text()) for li in element.find_all('li')]
                    if not items: continue
                    txBox = slide.shapes.add_textbox(left_margin + Inches(0.5), current_top, max_width - Inches(0.5), Inches(0.5 * len(items)))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    for i, item in enumerate(items):
                        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                        p.text = f"• {item}" if element.name == 'ul' else f"{i+1}. {item}"
                        p.font.size = Pt(16)
                        p.font.color.rgb = RGBColor(50, 50, 50)
                        p.level = 0
                    current_top += Inches(0.4 * len(items))
                    element.extract()
                    
                elif element.name == 'table':
                    rows = element.find_all('tr')
                    if not rows: continue
                    cols = max(len(row.find_all(['th', 'td'])) for row in rows)
                    if cols == 0: continue
                    
                    table_shape = slide.shapes.add_table(len(rows), cols, left_margin, current_top, max_width, Inches(0.4 * len(rows)))
                    table = table_shape.table
                    
                    for r_idx, row in enumerate(rows):
                        cells = row.find_all(['th', 'td'])
                        for c_idx, cell in enumerate(cells):
                            if c_idx < cols:
                                cell_frame = table.cell(r_idx, c_idx).text_frame
                                cell_frame.text = clean_text(cell.get_text())
                                for p in cell_frame.paragraphs:
                                    p.font.size = Pt(12)
                                    p.font.color.rgb = RGBColor(0, 0, 0)
                                    if cell.name == 'th':
                                        p.font.bold = True
                    
                    current_top += Inches(0.4 * len(rows) + 0.2)
                    element.extract()
                    
                elif element.name == 'img':
                    src = element.get('src')
                    if src and not src.startswith('http'):
                        img_path = os.path.join(base_dir, src.replace('/', os.sep))
                        if os.path.exists(img_path):
                            try:
                                slide.shapes.add_picture(img_path, left_margin, current_top, width=Inches(4))
                                current_top += Inches(3.2)
                            except Exception as e:
                                print(f"Failed to add image {img_path}: {e}")
                    element.extract()
                    
                elif element.name == 'div':
                    # Card handling
                    text = clean_text(element.get_text(separator=' | '))
                    if not text: continue
                    txBox = slide.shapes.add_textbox(left_margin, current_top, max_width, Inches(0.8))
                    txBox.fill.solid()
                    txBox.fill.fore_color.rgb = RGBColor(245, 245, 245)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = text
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    current_top += Inches(1.0)
                    element.extract()
                    
    prs.save(output_path)
    print(f"Saved presentation to {output_path}")

if __name__ == '__main__':
    base_dir = r"d:\DRIVE (Ai) Agents\00 Projects\Workplace CSR Slides"
    output_path = r"C:\Users\hashm\Desktop\PhD_Thesis_Presentation.pptx"
    process_html_to_pptx(base_dir, output_path)
